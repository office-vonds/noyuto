#!/usr/bin/env python3
"""Generate a premium 20-slide RIWAY presentation with luxury visual design."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Constants
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Premium Color Palette
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
MID_BG = RGBColor(0x1B, 0x3A, 0x5C)
CARD_BG = RGBColor(0x1E, 0x34, 0x54)
HIGHLIGHT_BG = RGBColor(0x16, 0x2D, 0x50)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WARM_GOLD = RGBColor(0xE8, 0xC5, 0x47)
LIGHT_GOLD = RGBColor(0xF0, 0xDF, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHITE = RGBColor(0xE8, 0xE8, 0xE8)
GRAY50 = RGBColor(0x80, 0x80, 0x80)
SHADOW_COLOR = RGBColor(0x08, 0x10, 0x1A)
BURGUNDY = RGBColor(0x8B, 0x1A, 0x1A)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(blank_layout)


def add_layered_bg(slide):
    """Create layered background for depth."""
    # Dark outer background
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid()
    s.fill.fore_color.rgb = DARK_BG
    s.line.fill.background()
    # Slightly lighter inner area (inset)
    inset = Inches(0.3)
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inset, inset,
                                 SLIDE_W - 2 * inset, SLIDE_H - 2 * inset)
    s2.fill.solid()
    s2.fill.fore_color.rgb = MID_BG
    s2.line.fill.background()


def add_decorative_circle(slide, left, top, size, color=LIGHT_GOLD):
    """Add a large semi-transparent decorative circle."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    # Simulate transparency with lighter color
    return c


def add_gold_bar_left(slide, top, height=Inches(0.8)):
    """Add thick gold vertical bar on left side of title."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.2), top, Pt(4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def add_thin_gold_line(slide, top, left=Inches(1.5), width=Inches(10.5)):
    """Add thin gold horizontal divider line."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY50
    p.alignment = PP_ALIGN.RIGHT


def add_title_with_bar(slide, text, top=Inches(0.6)):
    """Add slide title with gold vertical bar accent."""
    add_gold_bar_left(slide, top, Inches(0.7))
    txBox = slide.shapes.add_textbox(Inches(1.6), top, Inches(10.0), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_text(slide, text, left, top, width, height, size=Pt(22),
             color=SOFT_WHITE, align=PP_ALIGN.LEFT, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_card(slide, left, top, width, height, lines, gold_top=True,
             shadow=True, font_size=Pt(18), align=PP_ALIGN.LEFT):
    """Add a premium card with shadow and gold top accent."""
    # Shadow
    if shadow:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left + Inches(0.05), top + Inches(0.05), width, height)
        sh.fill.solid()
        sh.fill.fore_color.rgb = SHADOW_COLOR
        sh.line.fill.background()
    # Card body
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.fill.background()
    # Gold top accent bar
    if gold_top:
        gt = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
        gt.fill.solid()
        gt.fill.fore_color.rgb = GOLD
        gt.line.fill.background()
    # Text
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.25)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = SOFT_WHITE
        p.alignment = align
        p.space_after = Pt(4)
    return card


def add_number_circle(slide, left, top, number, size=Inches(0.5)):
    """Add gold circle with white number."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = GOLD
    c.line.fill.background()
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return c


def add_gold_diamond(slide, left, top, size=Inches(0.15)):
    """Add small gold diamond as bullet marker."""
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, size, size)
    d.fill.solid()
    d.fill.fore_color.rgb = GOLD
    d.line.fill.background()
    return d


def add_chevron(slide, left, top, width=Inches(0.4), height=Inches(0.5)):
    """Add gold chevron arrow."""
    ch = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    ch.fill.solid()
    ch.fill.fore_color.rgb = GOLD
    ch.line.fill.background()
    return ch


# ==================================================
# SLIDE 1 - Title
# ==================================================
slide = new_slide()
add_layered_bg(slide)
# Large decorative gold circle bleeding off bottom-right
add_decorative_circle(slide, Inches(9.5), Inches(4.5), Inches(5.0), LIGHT_GOLD)
add_text(slide, "あなたの人生を変える、たった一つの選択",
         Inches(1.5), Inches(2.0), Inches(8.0), Inches(1.5),
         size=Pt(48), color=WHITE, bold=True, align=PP_ALIGN.LEFT)
add_thin_gold_line(slide, Inches(3.5), Inches(1.5), Inches(6.0))
add_text(slide, "RIWAY Business Presentation",
         Inches(1.5), Inches(3.8), Inches(8.0), Inches(0.8),
         size=Pt(28), color=GOLD, align=PP_ALIGN.LEFT)
add_text(slide, "Confidential | For Invited Guests Only",
         Inches(1.5), Inches(6.5), Inches(6.0), Inches(0.5),
         size=Pt(12), color=GRAY50, align=PP_ALIGN.LEFT)
add_slide_number(slide, 1)

# ==================================================
# SLIDE 2 - Hook
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_decorative_circle(slide, Inches(-1.5), Inches(-1.5), Inches(4.0), LIGHT_GOLD)
add_text(slide, '日本人の2人に1人が、がんになる時代。',
         Inches(2.0), Inches(2.5), Inches(9.333), Inches(1.0),
         size=Pt(36), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 'あなたの健康対策は、本当に十分ですか？',
         Inches(2.0), Inches(3.8), Inches(9.333), Inches(1.0),
         size=Pt(36), color=GOLD, bold=True, align=PP_ALIGN.CENTER)
add_slide_number(slide, 2)

# ==================================================
# SLIDE 3 - Stats
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, '知っていますか？この現実を。')
add_thin_gold_line(slide, Inches(1.5), Inches(1.6), Inches(10.0))

stats = [
    ('男性9年・奣12年', '平均寿命と健康寿命の差'),
    ('約2,800万円', '生涯医療費'),
    ('わずか3%', '退職後の収入源を持つ人'),
]
card_w = Inches(3.2)
card_h = Inches(3.2)
card_top = Inches(2.5)
start_x = Inches(1.8)
gap = Inches(0.6)

for i, (number, label) in enumerate(stats):
    left = start_x + i * (card_w + gap)
    card = add_card(slide, left, card_top, card_w, card_h, ['', '', ''], font_size=Pt(18))
    # Number big and gold
    add_text(slide, number, left + Inches(0.3), card_top + Inches(0.5),
             card_w - Inches(0.6), Inches(1.2),
             size=Pt(54), color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    # Label below
    add_text(slide, label, left + Inches(0.3), card_top + Inches(2.0),
             card_w - Inches(0.6), Inches(0.8),
             size=Pt(18), color=SOFT_WHITE, align=PP_ALIGN.CENTER)
add_slide_number(slide, 3)

# ==================================================
# SLIDE 4 - Two fears
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, '現代人が抱える「2つの不安」')
add_thin_gold_line(slide, Inches(1.5), Inches(1.6), Inches(10.0))

col_w = Inches(5.0)
col_h = Inches(4.2)
col_top = Inches(2.3)

# Left card
card_l = add_card(slide, Inches(1.5), col_top, col_w, col_h,
                  ['', '', 'サプリは多すぎて選べない', '高額な健康法は続かない', '何が本物かわからない'],
                  font_size=Pt(20))
add_text(slide, '健康の不安', Inches(1.7), col_top + Inches(0.3),
         Inches(4.6), Inches(0.6), size=Pt(26), color=GOLD, bold=True)

# Right card
card_r = add_card(slide, Inches(7.0), col_top, col_w, col_h,
                  ['', '', '給料は上がらず物価は上がる', '年金だけでは不安', '副業したいが時間がない'],
                  font_size=Pt(20))
add_text(slide, '経済の不安', Inches(7.2), col_top + Inches(0.3),
         Inches(4.6), Inches(0.6), size=Pt(26), color=GOLD, bold=True)

# Add gold diamond bullets
for j in range(3):
    add_gold_diamond(slide, Inches(1.9), col_top + Inches(1.3) + j * Inches(0.65))
    add_gold_diamond(slide, Inches(7.4), col_top + Inches(1.3) + j * Inches(0.65))
add_slide_number(slide, 4)

# ==================================================
# SLIDE 5 - Barriers
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, 'なぜ、普通の選択肢では解決できないのか')
add_thin_gold_line(slide, Inches(1.5), Inches(1.6), Inches(10.0))

barriers = [
    (1, '時間の壁', '忙しくて健康管理ができない'),
    (2, 'コストの壁', '良いものほど高額で続かない'),
    (3, '情報の壁', '何が本物かわからない'),
]
for i, (num, title, desc) in enumerate(barriers):
    row_top = Inches(2.5) + i * Inches(1.5)
    add_card(slide, Inches(2.8), row_top, Inches(8.5), Inches(1.2),
             [title + '：' + desc], font_size=Pt(24), align=PP_ALIGN.LEFT)
    add_number_circle(slide, Inches(2.0), row_top + Inches(0.2), num, Inches(0.6))
add_slide_number(slide, 5)

# ==================================================
# SLIDE 6 - Turning point
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_decorative_circle(slide, Inches(-2.0), Inches(-2.0), Inches(5.0), LIGHT_GOLD)
add_decorative_circle(slide, Inches(10.5), Inches(5.0), Inches(4.0), LIGHT_GOLD)
add_text(slide, 'もし「健康」と「収入」を',
         Inches(2.0), Inches(2.2), Inches(9.333), Inches(1.0),
         size=Pt(40), color=GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, '同時に手に入れる方法があるとしたら？',
         Inches(2.0), Inches(3.5), Inches(9.333), Inches(1.0),
         size=Pt(40), color=GOLD, bold=True, align=PP_ALIGN.CENTER)
add_slide_number(slide, 6)

# ==================================================
# SLIDE 7 - RIWAY intro
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, 'RIWAY とは', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

riway_items = [
    '2008年シンガポール設立',
    'アジア15カ国以上に展開',
    'ブランドメッセージ「Live Young」',
    '唯一無二の製品力で急成長',
]
for i, item in enumerate(riway_items):
    y = Inches(2.0) + i * Inches(1.0)
    add_gold_diamond(slide, Inches(2.0), y + Inches(0.1))
    add_text(slide, item, Inches(2.4), y, Inches(9.0), Inches(0.6),
             size=Pt(24), color=SOFT_WHITE)
add_slide_number(slide, 7)

# ==================================================
# SLIDE 8 - Trust
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_decorative_circle(slide, Inches(10.0), Inches(-1.5), Inches(4.0), LIGHT_GOLD)
add_title_with_bar(slide, 'RIWAY の信頼性', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

trust_items = [
    'ニュージーランド自社工場',
    'GMP認証取得',
    '世界各国の政府認可',
    '15年以上の実績',
    '日本法人設立済み',
]
for i, item in enumerate(trust_items):
    y = Inches(2.0) + i * Inches(0.9)
    add_gold_diamond(slide, Inches(2.0), y + Inches(0.1))
    add_text(slide, item, Inches(2.4), y, Inches(9.0), Inches(0.6),
             size=Pt(24), color=SOFT_WHITE)
add_slide_number(slide, 8)

# ==================================================
# SLIDE 9 - Awards
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, '受賞・メディア実績', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

award_items = [
    'アジア太平洋トップ企業賞受賞',
    '各種ビジネスアワード多数',
    '世界15カ国以上で政府認可',
    '累計販売実績 数百万箱突破',
]
for i, item in enumerate(award_items):
    y = Inches(2.0) + i * Inches(1.0)
    add_gold_diamond(slide, Inches(2.0), y + Inches(0.1))
    add_text(slide, item, Inches(2.4), y, Inches(9.0), Inches(0.6),
             size=Pt(24), color=SOFT_WHITE)
add_slide_number(slide, 9)

# ==================================================
# SLIDE 10 - Product
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_decorative_circle(slide, Inches(9.0), Inches(3.5), Inches(5.0), LIGHT_GOLD)
add_text(slide, 'PURTIER PLACENTA',
         Inches(1.5), Inches(0.8), Inches(10.0), Inches(1.2),
         size=Pt(48), color=WHITE, bold=True, align=PP_ALIGN.LEFT)
add_text(slide, '世界唯一の「生きた細胞」カプセル技術',
         Inches(1.5), Inches(1.8), Inches(10.0), Inches(0.8),
         size=Pt(24), color=GOLD, align=PP_ALIGN.LEFT)
add_thin_gold_line(slide, Inches(2.5), Inches(1.5), Inches(8.0))

product_items = [
    'ニュージーランド産 鹿プラセンタ使用',
    '特許技術による生細胞カプセル化',
    '9つの厳選された天然成分を配合',
    '1箱60カプセル（1ヶ月分）',
]
for i, item in enumerate(product_items):
    y = Inches(3.0) + i * Inches(0.9)
    add_gold_diamond(slide, Inches(1.8), y + Inches(0.1))
    add_text(slide, item, Inches(2.2), y, Inches(8.0), Inches(0.6),
             size=Pt(24), color=SOFT_WHITE)
add_slide_number(slide, 10)

# ==================================================
# SLIDE 11 - Ingredients
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, 'PURTIER の主要成分', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

left_items = ['鹿プラセンタエキス', '海洋コラーゲンペプチド', '大豆イソフラボン', 'アロエベラエキス']
right_items = ['マリンコラーゲン', 'ボラージオイル', 'アビジン', 'リコピン', 'ルテイン']

# Left column cards
for i, item in enumerate(left_items):
    y = Inches(2.0) + i * Inches(1.05)
    add_card(slide, Inches(1.5), y, Inches(5.0), Inches(0.8), [item],
             font_size=Pt(22), align=PP_ALIGN.LEFT)
    add_gold_diamond(slide, Inches(1.8), y + Inches(0.3))

# Right column cards
for i, item in enumerate(right_items):
    y = Inches(2.0) + i * Inches(1.05)
    add_card(slide, Inches(7.0), y, Inches(5.0), Inches(0.8), [item],
             font_size=Pt(22), align=PP_ALIGN.LEFT)
    add_gold_diamond(slide, Inches(7.3), y + Inches(0.3))
add_slide_number(slide, 11)

# ==================================================
# SLIDE 12 - Testimonials
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, '愛用者の声', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

testimonials = [
    ('50代女性', '「3ヶ月で肌のハリが戻り、\n周りから若返ったと\n言われます」'),
    ('60代男性', '「毎朝の目覚めが劇的に\n変わりました。ゴルフの\nスコアも改善」'),
    ('40代女性', '「疲れにくくなり、仕事と\n育児の両立が楽に\nなりました」'),
]
card_w = Inches(3.4)
card_h = Inches(3.8)
start_x = Inches(1.5)
gap = Inches(0.5)

for i, (who, quote) in enumerate(testimonials):
    left = start_x + i * (card_w + gap)
    card_top = Inches(2.0)
    card = add_card(slide, left, card_top, card_w, card_h, ['', '', '', quote],
                    font_size=Pt(16), align=PP_ALIGN.CENTER)
    # Large gold quotation mark
    add_text(slide, '「', left + Inches(0.2), card_top + Inches(0.2),
             Inches(1.0), Inches(0.8), size=Pt(48), color=GOLD, bold=True)
    # Who label in gold
    add_text(slide, who, left + Inches(0.3), card_top + Inches(0.7),
             card_w - Inches(0.6), Inches(0.5), size=Pt(20), color=GOLD, bold=True,
             align=PP_ALIGN.CENTER)

add_text(slide, '※個人の感想です。効果には個人差があります。',
         Inches(1.0), Inches(6.3), Inches(11.333), Inches(0.5),
         size=Pt(12), color=GRAY50, align=PP_ALIGN.CENTER)
add_slide_number(slide, 12)

# ==================================================
# SLIDE 13 - Business model
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, '消費しながら収益を生む仕組み', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

steps = [
    '自分で使う → 健康を実感する',
    '良さを伝える → 共感が広がる',
    'チームで広がる → 収益が生まれる',
]
for i, text in enumerate(steps):
    row_top = Inches(2.3) + i * Inches(1.5)
    add_number_circle(slide, Inches(2.0), row_top + Inches(0.15), i + 1, Inches(0.7))
    add_card(slide, Inches(3.0), row_top, Inches(8.0), Inches(1.1),
             [text], font_size=Pt(24), align=PP_ALIGN.LEFT)
    # Chevron between steps
    if i < 2:
        add_chevron(slide, Inches(6.3), row_top + Inches(1.15), Inches(0.4), Inches(0.3))
add_slide_number(slide, 13)

# ==================================================
# SLIDE 14 - Compensation
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, '7つの報酬プラン', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

plans = [
    'リテールプロフィット',
    'パーソナルボーナス',
    'チームボーナス',
    'リーダーシップボーナス',
    'グローバルプールボーナス',
    'トラベルファンド',
    'カーファンド',
]
for i, plan in enumerate(plans):
    y = Inches(1.8) + i * Inches(0.7)
    # Alternating subtle background bars
    if i % 2 == 0:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(2.0), y, Inches(9.5), Inches(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = HIGHLIGHT_BG
        bar.line.fill.background()
    add_text(slide, str(i + 1) + '.', Inches(2.2), y + Inches(0.05),
             Inches(0.5), Inches(0.5), size=Pt(22), color=GOLD, bold=True)
    add_text(slide, plan, Inches(2.8), y + Inches(0.05),
             Inches(8.0), Inches(0.5), size=Pt(22), color=SOFT_WHITE)

add_text(slide, '※詳細は紹介者にお問い合わせください',
         Inches(1.0), Inches(6.5), Inches(11.333), Inches(0.5),
         size=Pt(12), color=GRAY50, align=PP_ALIGN.CENTER)
add_slide_number(slide, 14)

# ==================================================
# SLIDE 15 - Simulation
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, '収入シミュレーション', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

sims = [
    ('3人に紹介', '月5〜10万円'),
    ('9人のチーム', '月20〜50万円'),
    ('27人の組織', '月100万円以上'),
]
card_w = Inches(3.2)
card_h = Inches(3.0)
start_x = Inches(1.5)
gap = Inches(0.4)

for i, (label, amount) in enumerate(sims):
    left = start_x + i * (card_w + gap + Inches(0.3))
    card_top = Inches(2.2)
    add_card(slide, left, card_top, card_w, card_h, ['', '', ''],
             font_size=Pt(18))
    add_text(slide, label, left + Inches(0.2), card_top + Inches(0.4),
             card_w - Inches(0.4), Inches(0.6),
             size=Pt(20), color=SOFT_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, amount, left + Inches(0.2), card_top + Inches(1.2),
             card_w - Inches(0.4), Inches(1.0),
             size=Pt(42), color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    # Arrow between cards
    if i < 2:
        arrow_x = left + card_w + Inches(0.05)
        add_chevron(slide, arrow_x, card_top + Inches(1.2), Inches(0.35), Inches(0.5))

add_text(slide, '※シミュレーション例です。収入を保証するものではありません。',
         Inches(1.0), Inches(6.3), Inches(11.333), Inches(0.5),
         size=Pt(12), color=GRAY50, align=PP_ALIGN.CENTER)
add_slide_number(slide, 15)

# ==================================================
# SLIDE 16 - Timeline / Success
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, '成功者の軌跡', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

timeline = [
    ('1ヶ月目', '製品を体感、3人に紹介'),
    ('3ヶ月目', 'チーム9人、月収30万円達成'),
    ('6ヶ月目', '組織拡大、月収100万円突破'),
    ('12ヶ月目', '海外旅行ファンド獲得'),
]
# Vertical gold line on left
vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(2.8), Inches(2.0), Pt(2), Inches(4.5))
vline.fill.solid()
vline.fill.fore_color.rgb = GOLD
vline.line.fill.background()

for i, (period, desc) in enumerate(timeline):
    y = Inches(2.2) + i * Inches(1.1)
    # Gold dot on line
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(2.65), y + Inches(0.1), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GOLD
    dot.line.fill.background()
    # Period text
    add_text(slide, period, Inches(3.4), y,
             Inches(2.5), Inches(0.5), size=Pt(22), color=GOLD, bold=True)
    # Description
    add_text(slide, desc, Inches(5.5), y,
             Inches(6.0), Inches(0.5), size=Pt(22), color=SOFT_WHITE)

add_text(slide, '※個人の成果例です。成果には個人差があります。',
         Inches(1.0), Inches(6.8), Inches(11.333), Inches(0.5),
         size=Pt(12), color=GRAY50, align=PP_ALIGN.CENTER)
add_slide_number(slide, 16)

# ==================================================
# SLIDE 17 - Why now
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, 'なぜ「今」なのか', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

why_items = [
    ('★', '日本市場はまだ初期段階', '先行者利益が最大'),
    ('▲', '超高齢化社会', '健康需要は拡大し続ける'),
    ('●', '円安時代', '外貨建て収入の価値が高まる'),
]
card_w = Inches(3.3)
card_h = Inches(3.5)
start_x = Inches(1.5)
gap = Inches(0.5)

for i, (icon, title, desc) in enumerate(why_items):
    left = start_x + i * (card_w + gap)
    card_top = Inches(2.2)
    add_card(slide, left, card_top, card_w, card_h, ['', '', '', ''],
             font_size=Pt(18))
    # Icon shape (gold circle with symbol)
    add_text(slide, icon, left + Inches(1.2), card_top + Inches(0.3),
             Inches(1.0), Inches(0.8), size=Pt(36), color=GOLD, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, title, left + Inches(0.2), card_top + Inches(1.2),
             card_w - Inches(0.4), Inches(0.6), size=Pt(22), color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, desc, left + Inches(0.2), card_top + Inches(2.0),
             card_w - Inches(0.4), Inches(0.8), size=Pt(18), color=SOFT_WHITE,
             align=PP_ALIGN.CENTER)
add_slide_number(slide, 17)

# ==================================================
# SLIDE 18 - How to start
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, '始め方は、驚くほど簡単です', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

start_steps = [
    '今日、紹介者と話す',
    '製品を体感する',
    'ビジネスをスタート',
]
circle_size = Inches(1.0)
start_x = Inches(2.0)
gap_x = Inches(3.5)

for i, text in enumerate(start_steps):
    cx = start_x + i * gap_x
    cy = Inches(2.8)
    # Large gold numbered circle
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, circle_size, circle_size)
    c.fill.solid()
    c.fill.fore_color.rgb = GOLD
    c.line.fill.background()
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(36)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # Step text below
    add_text(slide, text, cx - Inches(0.5), cy + Inches(1.3),
             Inches(2.0), Inches(1.0), size=Pt(22), color=SOFT_WHITE,
             align=PP_ALIGN.CENTER)
add_slide_number(slide, 18)

# ==================================================
# SLIDE 19 - FAQ
# ==================================================
slide = new_slide()
add_layered_bg(slide)
add_title_with_bar(slide, 'よくある質問', top=Inches(0.5))
add_thin_gold_line(slide, Inches(1.4), Inches(1.6), Inches(10.0))

faqs = [
    ('Q: 営業経験がなくても大丈夫？', 'A: 充実したチームサポート体制があります'),
    ('Q: いくらから始められる？', 'A: 製品購入からスタート。詳細は紹介者へ'),
    ('Q: 本業と両立できる？', 'A: スキマ時間で取り組めます'),
]
for i, (q, a) in enumerate(faqs):
    y = Inches(1.8) + i * Inches(1.7)
    card = add_card(slide, Inches(2.0), y, Inches(9.5), Inches(1.4),
                    [q, a], font_size=Pt(20), align=PP_ALIGN.LEFT)
    tf = card.text_frame
    for p in tf.paragraphs:
        if p.text.startswith('Q:'):
            p.font.color.rgb = GOLD
            p.font.bold = True
            p.font.size = Pt(22)
add_slide_number(slide, 19)

# ==================================================
# SLIDE 20 - CTA
# ==================================================
slide = new_slide()
add_layered_bg(slide)
# Dramatic decorative elements
add_decorative_circle(slide, Inches(-2.0), Inches(3.5), Inches(5.0), LIGHT_GOLD)
add_decorative_circle(slide, Inches(10.0), Inches(-2.0), Inches(5.0), LIGHT_GOLD)
# Top thin gold lines
add_thin_gold_line(slide, Inches(2.5), Inches(3.0), Inches(7.333))
add_thin_gold_line(slide, Inches(5.0), Inches(3.0), Inches(7.333))

add_text(slide, 'あなたの人生を変える一歩は、',
         Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.0),
         size=Pt(38), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, '今日ここから始まります。',
         Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.0),
         size=Pt(38), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Gold CTA box
cta_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(3.0), Inches(4.0), Inches(7.333), Inches(1.2))
cta_box.fill.solid()
cta_box.fill.fore_color.rgb = GOLD
cta_box.line.fill.background()
tf = cta_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '今すぐ、紹介者にご連絡ください。'
p.font.size = Pt(30)
p.font.color.rgb = DARK_BG
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

add_text(slide, '本日の資料に関するお問い合わせは、お近くのスタッフまで',
         Inches(1.0), Inches(6.3), Inches(11.333), Inches(0.6),
         size=Pt(14), color=GRAY50, align=PP_ALIGN.CENTER)
add_slide_number(slide, 20)

# ==================================================
# Save
# ==================================================
output_path = "/home/user/noyuto/RIWAY_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
